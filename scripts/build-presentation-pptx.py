# -*- coding: utf-8 -*-
"""Build Idolelic pitch deck as PowerPoint for Google Drive / Slides."""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt, Emu

ROOT = Path(__file__).resolve().parents[1]
ASSETS = ROOT / "docs" / "presentation-assets"
OUT = ROOT / "docs" / "presentation-idolelic.pptx"

PINK = RGBColor(0xE8, 0x3E, 0x8C)
PINK_DEEP = RGBColor(0xC2, 0x18, 0x5B)
TEXT = RGBColor(0x1A, 0x1A, 0x1A)
MUTED = RGBColor(0x5C, 0x5C, 0x5C)
OK = RGBColor(0x1F, 0x9D, 0x6A)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
SOFT = RGBColor(0xFF, 0xF0, 0xF6)
CARD = RGBColor(0xF7, 0xF7, 0xF8)


def set_run(run, *, size=14, bold=False, color=TEXT):
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = "Yu Gothic UI"


def add_textbox(slide, left, top, width, height):
    return slide.shapes.add_textbox(left, top, width, height)


def write_lines(tf, lines, *, size=14, color=TEXT, bold=False, space_after=8):
    tf.clear()
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(space_after)
        run = p.add_run()
        run.text = line
        set_run(run, size=size, bold=bold, color=color)


def add_section_label(slide, text):
    box = add_textbox(slide, Inches(0.55), Inches(0.28), Inches(9), Inches(0.35))
    write_lines(box.text_frame, [text], size=12, bold=True, color=PINK, space_after=0)


def add_title(slide, text):
    box = add_textbox(slide, Inches(0.55), Inches(0.55), Inches(9), Inches(0.85))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    set_run(run, size=26, bold=True, color=TEXT)


def add_bullets(slide, items, left, top, width, height):
    box = add_textbox(slide, left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    for i, (head, body) in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(14)
        r1 = p.add_run()
        r1.text = f"✓ {head}"
        set_run(r1, size=14, bold=True, color=TEXT)
        p2 = tf.add_paragraph()
        p2.space_after = Pt(10)
        r2 = p2.add_run()
        r2.text = body
        set_run(r2, size=12, color=MUTED)


def add_tags(slide, tags, top):
    from pptx.enum.shapes import MSO_SHAPE

    left = Inches(0.55)
    for tag, ok in tags:
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            left,
            top,
            Inches(max(1.4, 0.22 * len(tag))),
            Inches(0.34),
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(0xE8, 0xF8, 0xF0) if ok else SOFT
        shape.line.fill.background()
        tf = shape.text_frame
        tf.word_wrap = False
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = tag
        set_run(run, size=10, bold=True, color=OK if ok else PINK_DEEP)
        left += shape.width + Inches(0.12)


def add_image_if_exists(slide, name, left, top, width):
    path = ASSETS / name
    if not path.exists():
        return
    slide.shapes.add_picture(str(path), left, top, width=width)


def add_footer(slide, page, total=6):
    box = add_textbox(slide, Inches(0.55), Inches(6.95), Inches(8), Inches(0.3))
    write_lines(
        box.text_frame,
        [f"Idolelic ピッチプレゼンテーション　　{page} / {total}"],
        size=10,
        color=MUTED,
        space_after=0,
    )


def card(slide, left, top, width, height, title, lines, accent=PINK):
    from pptx.enum.shapes import MSO_SHAPE

    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = CARD
    shape.line.color.rgb = RGBColor(0xEE, 0xEE, 0xEE)
    title_box = add_textbox(slide, left + Inches(0.15), top + Inches(0.12), width - Inches(0.3), Inches(0.35))
    write_lines(title_box.text_frame, [title], size=13, bold=True, color=accent, space_after=0)
    body = add_textbox(slide, left + Inches(0.15), top + Inches(0.45), width - Inches(0.3), height - Inches(0.55))
    write_lines(body.text_frame, lines, size=11, color=TEXT, space_after=4)


def build():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    # --- Slide 1 ---
    s = prs.slides.add_slide(blank)
    add_section_label(s, "01 / INTRODUCTION")
    add_title(s, "聖地巡礼をしながら健康になれる推し活アプリ")
    add_bullets(
        s,
        [
            ("推しへの熱量を歩行へ：", "楽しんでいるうちに自然と歩数が伸び、運動不足を解消する聖地巡礼マップアプリ。"),
            ("コンセプトの進化：", "「今推し」だけでなく、時代を築いたアイドルの聖地を忘れないアーカイブとしても設計。"),
            ("いま動いていること：", "公式聖地65件の地図、MV埋め込み再生、アプリ内ナビ＋歩数、Supabase認証・オーナー管理まで実装済み。"),
        ],
        Inches(0.55),
        Inches(1.55),
        Inches(7.2),
        Inches(3.8),
    )
    add_tags(
        s,
        [("アプリ名: Idolelic", False), ("旧称 / リポジトリ: osimap", False), ("MVP ほぼ完成", True)],
        Inches(5.5),
    )
    add_image_if_exists(s, "home-map.png", Inches(8.3), Inches(1.5), Inches(4.4))
    cap = add_textbox(s, Inches(8.3), Inches(6.35), Inches(4.4), Inches(0.3))
    write_lines(cap.text_frame, ["実機画面：聖地マップ（公式65件）"], size=10, color=MUTED, space_after=0)
    add_footer(s, 1)

    # --- Slide 2 ---
    s = prs.slides.add_slide(blank)
    add_section_label(s, "02 / PROBLEM & TARGET")
    add_title(s, "オタクの不健康問題と、46兆円の壁")
    add_bullets(
        s,
        [
            ("深刻なオタクの運動不足：", "自宅での配信視聴など推し活の室内化により、生活習慣病リスクが高まっている。"),
            ("社会課題としての医療費：", "日本の年間国民医療費は約46兆円規模。予防医療・行動変容が急務。"),
            ("巡礼の手間：", "毎回サイトを横断し Google Maps に手作業でピン留めするハードルが高い。聖地情報が散逸しやすい。"),
        ],
        Inches(0.55),
        Inches(1.55),
        Inches(7.0),
        Inches(3.6),
    )
    add_tags(
        s,
        [("⚠ 医療費抑制へのアプローチ", False), ("聖地アーカイブ × 歩行習慣", False)],
        Inches(5.5),
    )
    from pptx.enum.shapes import MSO_SHAPE

    callout = s.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.0), Inches(1.55), Inches(4.7), Inches(4.2)
    )
    callout.fill.solid()
    callout.fill.fore_color.rgb = SOFT
    callout.line.fill.background()
    big = add_textbox(s, Inches(8.25), Inches(1.85), Inches(4.2), Inches(1.0))
    write_lines(big.text_frame, ["46兆円"], size=40, bold=True, color=PINK, space_after=0)
    sub = add_textbox(s, Inches(8.25), Inches(2.85), Inches(4.2), Inches(0.4))
    write_lines(sub.text_frame, ["日本の年間国民医療費（規模感）"], size=12, bold=True, color=PINK_DEEP, space_after=0)
    body = add_textbox(s, Inches(8.25), Inches(3.4), Inches(4.2), Inches(2.0))
    write_lines(
        body.text_frame,
        [
            "「推しに会いに行く」「お気に入りの場所を歩く」というオタク心理からアプローチする、画期的な行動変容。",
            "Idolelic は“歩く理由”を聖地巡礼に置き、継続しやすい健康習慣をつくる。",
        ],
        size=12,
        color=TEXT,
        space_after=8,
    )
    add_footer(s, 2)

    # --- Slide 3 ---
    s = prs.slides.add_slide(blank)
    add_section_label(s, "03 / CORE TECHNOLOGY")
    add_title(s, "不正を排除する、判定ロジックの信頼性")
    add_bullets(
        s,
        [
            ("時速10km/h未満の制限：", "自転車・電車などの乗り物移動を検知し、有効歩数への加算をブロック。"),
            ("5区間移動平均の実装：", "GPSのジャンプノイズを除外し、安定した徒歩距離のみを計測。"),
            ("歩数×GPS整合チェック：", "加速度センサーの生歩数とGPS移動距離を照合。振り子疑い・乗り物時は除外／上限キャップ。"),
        ],
        Inches(0.55),
        Inches(1.55),
        Inches(7.0),
        Inches(3.6),
    )
    add_tags(
        s,
        [("単体テスト通過（Vitest）", True), ("徒歩判定バリデータ", False), ("実機 HTTPS 検証済み", False)],
        Inches(5.5),
    )
    lab = s.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.0), Inches(1.55), Inches(4.7), Inches(4.0)
    )
    lab.fill.solid()
    lab.fill.fore_color.rgb = RGBColor(0x1A, 0x1A, 0x1A)
    lab.line.fill.background()
    lab_title = add_textbox(s, Inches(8.25), Inches(1.75), Inches(4.2), Inches(0.4))
    write_lines(lab_title.text_frame, ["GPS Lab / 速度フィルタ"], size=14, bold=True, color=WHITE, space_after=0)
    lab_body = add_textbox(s, Inches(8.25), Inches(2.35), Inches(4.2), Inches(3.0))
    write_lines(
        lab_body.text_frame,
        [
            "閾値　　10 km/h 未満 = 徒歩",
            "判定　　徒歩（有効）",
            "瞬間速度　0.0 km/h",
            "移動平均　0.0 km/h",
            "生歩数 / 有効歩数　照合中",
            "状態　　データ収集中…",
        ],
        size=13,
        color=WHITE,
        space_after=10,
    )
    add_footer(s, 3)

    # --- Slide 4 ---
    s = prs.slides.add_slide(blank)
    add_section_label(s, "04 / APPLICATION FEATURES")
    add_title(s, "オタクに寄り添う、実用的なUX設計")
    add_bullets(
        s,
        [
            ("実データ聖地マップ：", "12グループ・公式65件の MV ロケ地を登録。グループ／カテゴリ／地方フィルタ、YouTube MV 埋め込み再生。"),
            ("巡礼ナビと歩数：", "アプリ内徒歩ナビ（OSRM）と Google Maps 連携。GPS×加速度の有効歩数計測で歩きながら記録。"),
            ("みんなで残す基盤：", "住所→座標のジオコーディング、Supabase 認証・聖地永続化、オーナーによるデータ管理画面。"),
        ],
        Inches(0.55),
        Inches(1.55),
        Inches(7.0),
        Inches(3.5),
    )
    add_tags(
        s,
        [("MV埋め込み", False), ("アプリ内ナビ", False), ("公式65件シード済", True)],
        Inches(5.5),
    )
    add_image_if_exists(s, "spot-detail.png", Inches(8.1), Inches(1.45), Inches(2.35))
    add_image_if_exists(s, "register.png", Inches(10.6), Inches(1.45), Inches(2.35))
    c1 = add_textbox(s, Inches(8.1), Inches(5.9), Inches(2.35), Inches(0.3))
    write_lines(c1.text_frame, ["聖地詳細・案内"], size=10, color=MUTED, space_after=0)
    c2 = add_textbox(s, Inches(10.6), Inches(5.9), Inches(2.35), Inches(0.3))
    write_lines(c2.text_frame, ["聖地登録"], size=10, color=MUTED, space_after=0)
    add_footer(s, 4)

    # --- Slide 5 ---
    s = prs.slides.add_slide(blank)
    add_section_label(s, "05 / PROGRESS & TECH")
    add_title(s, "現在の進捗と、使用技術")
    card(
        s,
        Inches(0.55),
        Inches(1.5),
        Inches(2.7),
        Inches(2.5),
        "完了",
        [
            "・GPS・歩数コア／実機検証",
            "・地図・詳細・MV埋め込み",
            "・アプリ内ナビ／聖地登録",
            "・Auth・公式65件・管理画面",
        ],
        accent=OK,
    )
    card(
        s,
        Inches(3.4),
        Inches(1.5),
        Inches(2.7),
        Inches(2.5),
        "直前",
        [
            "・Vercel 本番デプロイ",
            "・本番 URL の Auth 設定",
            "・通しスモークテスト",
        ],
        accent=PINK,
    )
    card(
        s,
        Inches(6.25),
        Inches(1.5),
        Inches(2.7),
        Inches(2.5),
        "夏休み",
        [
            "・一般公開・運用開始",
            "・ユーザー獲得",
            "・掲示板DBなど拡張",
        ],
        accent=PINK_DEEP,
    )
    card(
        s,
        Inches(0.55),
        Inches(4.2),
        Inches(4.1),
        Inches(2.0),
        "Tech Stack",
        [
            "Next.js 16 / React 19 — App Router・API",
            "Supabase / Leaflet — Auth・聖地DB・地図",
            "Geolocation / 加速度 — 歩数照合",
            "OSRM / YouTube / Vitest — ナビ・MV・テスト",
        ],
    )
    add_image_if_exists(s, "mypage.png", Inches(9.1), Inches(1.5), Inches(3.6))
    cap = add_textbox(s, Inches(9.1), Inches(6.35), Inches(3.6), Inches(0.3))
    write_lines(cap.text_frame, ["実機画面：マイページ"], size=10, color=MUTED, space_after=0)
    add_footer(s, 5)

    # --- Slide 6 ---
    s = prs.slides.add_slide(blank)
    add_section_label(s, "06 / PERFORMANCE & ROADMAP")
    add_title(s, "MVPは実装完了目前。次は本番公開でトラクション獲得")
    milestones = [
        ("1. 機能追加", "コア歩数計＆マップの基本実装完了", True),
        ("2. 実機検証", "HTTPSでGPS/歩数ログ確認", True),
        ("3. MVP実装", "聖地65件・MV・ナビ・Auth・管理", True),
        ("4. 本番公開", "Vercelデプロイと通し確認", False),
        ("5. 夏休み", "トラクション獲得・掲示板DB等", False),
    ]
    x = Inches(0.45)
    for title, desc, done in milestones:
        card(
            s,
            x,
            Inches(1.7),
            Inches(2.35),
            Inches(3.2),
            title,
            [desc],
            accent=OK if done else PINK,
        )
        x += Inches(2.5)
    add_tags(
        s,
        [("コア〜MVP機能 完了", True), ("いま: Vercel公開", False), ("夏休み: トラクション獲得", False)],
        Inches(5.4),
    )
    demo = add_textbox(s, Inches(0.55), Inches(6.2), Inches(12), Inches(0.4))
    write_lines(
        demo.text_frame,
        ["デモURL: https://idolelic.vercel.app/home"],
        size=14,
        bold=True,
        color=PINK_DEEP,
        space_after=0,
    )
    add_footer(s, 6)

    prs.save(OUT)
    print(f"Wrote {OUT}")


if __name__ == "__main__":
    build()
